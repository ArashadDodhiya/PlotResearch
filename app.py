import os
import json
import uuid
from flask import Flask, render_template, request, jsonify, redirect, url_for, send_file
from werkzeug.utils import secure_filename
from config import Config
from analyzer import DocumentAnalyzer
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO

app = Flask(__name__)
app.config.from_object(Config)
Config.init_app(app)

# Initialize analyzer
analyzer = DocumentAnalyzer()

def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/')
def index():
    """Landing page with upload interface"""
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file upload and analysis"""
    
    # Check if file is present
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type. Allowed: PDF, DOCX, TXT'}), 400
    
    try:
        # Save file
        filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
        file.save(file_path)
        
        # Get options
        custom_prompt = request.form.get('custom_prompt', '')
        summary_length = request.form.get('summary_length', 'Standard')
        
        # Analyze document
        dashboard_data = analyzer.analyze_document(file_path, custom_prompt, summary_length)
        
        # Save extracted text for chat context
        if 'text_content' in dashboard_data:
            text_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{unique_id}.txt")
            with open(text_path, 'w', encoding='utf-8') as f:
                f.write(dashboard_data['text_content'])
        
        # Check if it's custom HTML or JSON format
        if dashboard_data.get('is_custom_html'):
            # Inject Chat Widget
            html_content = dashboard_data['html_content']
            chat_widget_code = f"""
            <!-- Chat Widget -->
            <div id="chat-widget" style="position: fixed; bottom: 20px; right: 20px; z-index: 9999; font-family: system-ui, -apple-system, sans-serif;">
                <button id="chat-toggle" onclick="toggleChat()" style="background: #000; color: white; border: none; padding: 15px; border-radius: 50%; cursor: pointer; box-shadow: 0 4px 12px rgba(0,0,0,0.15); width: 60px; height: 60px; display: flex; align-items: center; justify-content: center; transition: transform 0.2s;">
                    <svg width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 15a2 2 0 0 1-2 2H7l-4 4V5a2 2 0 0 1 2-2h14a2 2 0 0 1 2 2z"></path></svg>
                </button>
                <div id="chat-window" style="display: none; position: absolute; bottom: 80px; right: 0; width: 350px; height: 500px; background: white; border-radius: 12px; box-shadow: 0 5px 20px rgba(0,0,0,0.2); flex-direction: column; overflow: hidden; border: 1px solid #eee;">
                    <div style="padding: 15px; background: #f8f9fa; border-bottom: 1px solid #eee; display: flex; justify-content: space-between; align-items: center;">
                        <h3 style="margin: 0; font-size: 16px; font-weight: 600;">Chat with Document</h3>
                        <button onclick="toggleChat()" style="background: none; border: none; cursor: pointer; font-size: 18px;">&times;</button>
                    </div>
                    <div id="chat-messages" style="flex: 1; padding: 15px; overflow-y: auto; background: #fff;">
                        <div style="margin-bottom: 10px; color: #666; font-size: 14px;">Ask me anything about this document!</div>
                    </div>
                    <div style="padding: 15px; border-top: 1px solid #eee; display: flex; gap: 10px;">
                        <input type="text" id="chat-input" placeholder="Type a question..." style="flex: 1; padding: 10px; border: 1px solid #ddd; border-radius: 6px; outline: none;" onkeypress="handleKeyPress(event)">
                        <button onclick="sendMessage()" style="background: #000; color: white; border: none; padding: 10px 15px; border-radius: 6px; cursor: pointer;">Send</button>
                    </div>
                </div>
            </div>
            <script>
                const DASHBOARD_ID = "{unique_id}";
                function toggleChat() {{
                    const w = document.getElementById('chat-window');
                    const b = document.getElementById('chat-toggle');
                    if (w.style.display === 'none') {{
                        w.style.display = 'flex';
                        b.style.transform = 'scale(0.9)';
                    }} else {{
                        w.style.display = 'none';
                        b.style.transform = 'scale(1)';
                    }}
                }}
                function handleKeyPress(e) {{
                    if (e.key === 'Enter') sendMessage();
                }}
                async function sendMessage() {{
                    const input = document.getElementById('chat-input');
                    const msgs = document.getElementById('chat-messages');
                    const text = input.value.trim();
                    if (!text) return;
                    
                    // Add user message
                    msgs.innerHTML += `<div style="margin: 10px 0; text-align: right;"><span style="background: #000; color: white; padding: 8px 12px; border-radius: 12px 12px 0 12px; display: inline-block; font-size: 14px;">${{text}}</span></div>`;
                    input.value = '';
                    msgs.scrollTop = msgs.scrollHeight;
                    
                    // Add loading
                    const loadingId = 'loading-' + Date.now();
                    msgs.innerHTML += `<div id="${{loadingId}}" style="margin: 10px 0; text-align: left;"><span style="background: #f1f1f1; color: #333; padding: 8px 12px; border-radius: 12px 12px 12px 0; display: inline-block; font-size: 14px;">Thinking...</span></div>`;
                    msgs.scrollTop = msgs.scrollHeight;
                    
                    try {{
                        const res = await fetch('/api/chat', {{
                            method: 'POST',
                            headers: {{ 'Content-Type': 'application/json' }},
                            body: JSON.stringify({{ dashboard_id: DASHBOARD_ID, message: text }})
                        }});
                        const data = await res.json();
                        
                        // Remove loading and add response
                        document.getElementById(loadingId).remove();
                        msgs.innerHTML += `<div style="margin: 10px 0; text-align: left;"><span style="background: #f1f1f1; color: #333; padding: 8px 12px; border-radius: 12px 12px 12px 0; display: inline-block; font-size: 14px;">${{data.response}}</span></div>`;
                    }} catch (e) {{
                        document.getElementById(loadingId).remove();
                        msgs.innerHTML += `<div style="margin: 10px 0; text-align: left; color: red;">Error: ${{e.message}}</div>`;
                    }}
                    msgs.scrollTop = msgs.scrollHeight;
                }}
            </script>
            """
            
            # Insert before </body>
            if '</body>' in html_content:
                html_content = html_content.replace('</body>', f'{chat_widget_code}</body>')
            else:
                html_content += chat_widget_code
            
            # Save HTML directly
            html_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{unique_id}.html")
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
        else:
            # Save dashboard JSON (fallback)
            dashboard_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{unique_id}.json")
            with open(dashboard_path, 'w', encoding='utf-8') as f:
                json.dump(dashboard_data, f, indent=2)
        
        # Clean up uploaded file
        os.remove(file_path)
        
        return jsonify({
            'success': True,
            'dashboard_id': unique_id,
            'is_custom_html': dashboard_data.get('is_custom_html', False),
            'redirect': url_for('dashboard', dashboard_id=unique_id)
        })
    
    except Exception as e:
        return jsonify({'error': f'Analysis failed: {str(e)}'}), 500

@app.route('/api/chat', methods=['POST'])
def chat_api():
    """Handle chat requests"""
    data = request.json
    dashboard_id = data.get('dashboard_id')
    message = data.get('message')
    
    if not dashboard_id or not message:
        return jsonify({'error': 'Missing data'}), 400
        
    # Load document text
    text_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.txt")
    if not os.path.exists(text_path):
        return jsonify({'error': 'Document context not found'}), 404
        
    with open(text_path, 'r', encoding='utf-8') as f:
        text = f.read()
        
    # Generate response
    response = analyzer.chat_with_document(text, message)
    return jsonify({'response': response})

@app.route('/dashboard/<dashboard_id>')
def dashboard(dashboard_id):
    """Display dashboard"""
    
    # Check for custom HTML first
    html_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.html")
    if os.path.exists(html_path):
        # Serve custom HTML directly
        with open(html_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    # Fallback to JSON template
    dashboard_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.json")
    
    if not os.path.exists(dashboard_path):
        return "Dashboard not found", 404
    
    with open(dashboard_path, 'r', encoding='utf-8') as f:
        dashboard_data = json.load(f)
    
    return render_template('dashboard.html', 
                         dashboard=dashboard_data, 
                         dashboard_id=dashboard_id)

@app.route('/export/<dashboard_id>')
def export_dashboard(dashboard_id):
    """Export dashboard as standalone HTML"""
    
    # Check for custom HTML first
    html_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.html")
    if os.path.exists(html_path):
        # Serve custom HTML directly
        with open(html_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    # Fallback to JSON template
    dashboard_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.json")
    
    if not os.path.exists(dashboard_path):
        return "Dashboard not found", 404
    
    with open(dashboard_path, 'r', encoding='utf-8') as f:
        dashboard_data = json.load(f)
    
    return render_template('export.html', dashboard=dashboard_data)

@app.route('/api/dashboard/<dashboard_id>')
def get_dashboard_json(dashboard_id):
    """Get dashboard JSON data"""
    
    dashboard_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.json")
    
    if not os.path.exists(dashboard_path):
        return jsonify({'error': 'Dashboard not found'}), 404
    
    with open(dashboard_path, 'r', encoding='utf-8') as f:
        dashboard_data = json.load(f)
    
    return jsonify(dashboard_data)

@app.route('/export/pptx/<dashboard_id>')
def export_pptx(dashboard_id):
    """Export dashboard as PowerPoint presentation"""
    
    # Check for custom HTML first
    html_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.html")
    text_path = os.path.join(app.config['DASHBOARD_FOLDER'], f"{dashboard_id}.txt")
    
    if not os.path.exists(text_path):
        return "Dashboard not found", 404
    
    # Read document text
    with open(text_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Extract title from text (first line or first 100 chars)
    doc_title = text.split('\n')[0][:100] if text else "Research Dashboard"
    title.text = doc_title
    subtitle.text = "Generated Dashboard Report"
    
    # Content Slide - Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    
    tf = body_shape.text_frame
    tf.text = "Key Highlights:"
    
    # Extract first few paragraphs
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()][:5]
    for para in paragraphs[:3]:
        p = tf.add_paragraph()
        p.text = para[:200] + "..." if len(para) > 200 else para
        p.level = 1
    
    # Content Slide - Key Findings
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Key Findings"
    
    tf = body_shape.text_frame
    tf.text = "Main Insights:"
    
    # Add more content
    for para in paragraphs[3:6]:
        p = tf.add_paragraph()
        p.text = para[:200] + "..." if len(para) > 200 else para
        p.level = 1
    
    # Save to BytesIO
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return send_file(
        pptx_io,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name=f'dashboard_{dashboard_id}.pptx'
    )

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
